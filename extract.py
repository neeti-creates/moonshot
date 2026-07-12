"""
MoonshotHunt — submission asset extraction.

Each extractor is independently wrapped in try/except so a single bad source
never kills the whole submission. `build_context()` fans out over uploaded
files + URLs and concatenates labelled source blocks, truncating only to
protect the free-tier model from runaway prompts.

Extractors accept a *filepath* (str). For Flask uploads, save the FileStorage
to a temp file first and pass that path (see app.py). `build_context` takes
lightweight source objects: {filename, filepath} for files, and plain URL
strings for links.
"""
import os
import logging
from typing import List, Dict, Any

log = logging.getLogger("moonshot.extract")


# --------------------------------------------------------------------------
# PDF (text-based, with OCR fallback for scanned/image-heavy decks)
# --------------------------------------------------------------------------
def extract_pdf(filepath: str) -> Dict[str, Any]:
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        text = text.strip()
        if len(text) < 50:  # likely scanned / image-based
            log.info("PDF text < 50 chars, falling back to OCR: %s", filepath)
            text = ocr_pdf(filepath)
        return {"ok": True, "text": text}
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "error": str(e)}


def ocr_pdf(filepath: str) -> str:
    try:
        from pdf2image import convert_from_path
        import pytesseract
        images = convert_from_path(filepath)
        text = ""
        for i, img in enumerate(images):
            text += f"[Slide {i + 1}]\n" + pytesseract.image_to_string(img) + "\n"
        return text.strip()
    except Exception as e:  # noqa: BLE001
        # OCR path failed (missing tesseract/poppler or corrupt file) — degrade gracefully
        log.warning("OCR fallback failed for %s: %s", filepath, e)
        return ""


# --------------------------------------------------------------------------
# PPTX (text frames + speaker notes)
# --------------------------------------------------------------------------
def extract_pptx(filepath: str) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"[Slide {i + 1}]\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                text += f"[Notes: {slide.notes_slide.notes_text_frame.text}]\n"
        return {"ok": True, "text": text.strip()}
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "error": str(e)}


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------
def extract_docx(filepath: str) -> Dict[str, Any]:
    try:
        import docx
        d = docx.Document(filepath)
        text = "\n".join(p.text for p in d.paragraphs)
        return {"ok": True, "text": text.strip()}
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "error": str(e)}


# --------------------------------------------------------------------------
# URL (website / article)
# --------------------------------------------------------------------------
def extract_url(url: str, timeout: int = 10) -> Dict[str, Any]:
    try:
        import requests
        from bs4 import BeautifulSoup
        headers = {"User-Agent": "Mozilla/5.0 (compatible; MoonshotHuntBot/1.0)"}
        resp = requests.get(url, headers=headers, timeout=timeout)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # collapse long blank runs
        lines = [ln for ln in (l.strip() for l in text.splitlines()) if ln]
        text = "\n".join(lines)
        return {"ok": True, "text": text[:8000]}  # cap to avoid runaway pages
    except Exception as e:  # noqa: BLE001
        return {"ok": False, "error": str(e)}


# --------------------------------------------------------------------------
# Orchestration
# --------------------------------------------------------------------------
_EXTRACTORS = {"pdf": extract_pdf, "pptx": extract_pptx, "docx": extract_docx}


def build_context(files: List[Dict[str, str]], urls: List[str]) -> str:
    """Concatenate all sources into one labelled blob.

    `files` is a list of {"filename": str, "filepath": str}.
    `urls` is a list of URL strings. Failures are logged and skipped.
    """
    sources: List[str] = []

    for f in files or []:
        filename = f.get("filename", "")
        filepath = f.get("filepath", "")
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        fn = _EXTRACTORS.get(ext)
        if not fn or not filepath:
            log.warning("Skipping unsupported/unreadable file: %s", filename)
            continue
        result = fn(filepath)
        if result["ok"]:
            sources.append(f"[FROM: {filename}]\n{result['text']}")
        else:
            log.warning("Extraction failed for %s: %s", filename, result["error"])

    for url in urls or []:
        result = extract_url(url)
        if result["ok"]:
            sources.append(f"[FROM: {url}]\n{result['text']}")
        else:
            log.warning("Extraction failed for %s: %s", url, result["error"])

    combined = "\n\n---\n\n".join(sources)
    # truncate if huge, to protect free-tier model from timeout
    if len(combined) > 15000:
        combined = combined[:15000] + "\n\n[content truncated for length]"
    return combined
